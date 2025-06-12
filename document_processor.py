import streamlit as st
import io
import pandas as pd
from pypdf import PdfReader
from pptx import Presentation
import docx
import xlrd # For .xls files
from bs4 import BeautifulSoup

from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS

# --- File Text Extraction Functions ---
def extract_text_from_pdf(pdf_file_content: io.BytesIO):
    """Extracts text from a PDF file."""
    text = ""
    try:
        pdf_reader = PdfReader(pdf_file_content)
        for page in pdf_reader.pages:
            text += page.extract_text() or "" # Ensure text is not None
    except Exception as e:
        st.error(f"Error processing PDF: {e}")
    return text

def extract_text_from_ppt(ppt_file_content: io.BytesIO):
    """Extracts text from a PPTX file."""
    text = ""
    try:
        presentation = Presentation(ppt_file_content)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        st.error(f"Error processing PPTX: {e}")
    return text

def extract_text_from_py(py_file_content: io.BytesIO):
    """Extracts text from a Python (.py) file."""
    text = ""
    try:
        text = py_file_content.read().decode("utf-8")
    except Exception as e:
        st.error(f"Error processing Python file: {e}")
    return text

def extract_text_from_docx(docx_file_content: io.BytesIO):
    """Extracts text from a DOCX file."""
    text = ""
    try:
        doc = docx.Document(docx_file_content)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    except Exception as e:
        st.error(f"Error processing DOCX: {e}")
    return text

def extract_text_from_excel(excel_file_content: io.BytesIO):
    """Extracts text from an Excel (.xls or .xlsx) file."""
    text = ""
    try:
        # Try openpyxl for .xlsx first
        try:
            df = pd.read_excel(excel_file_content, engine='openpyxl')
            text = df.to_string(index=False)
        except Exception:
            # Fallback to xlrd for .xls (older formats)
            excel_file_content.seek(0) # Reset stream position
            workbook = xlrd.open_workbook(file_contents=excel_file_content.read())
            for sheet in workbook.sheets():
                for row in range(sheet.nrows):
                    for col in range(sheet.ncols):
                        cell_value = sheet.cell(row, col).value
                        if isinstance(cell_value, str):
                            text += cell_value + " "
    except Exception as e:
        st.error(f"Error processing Excel file: {e}")
    return text

def extract_text_from_csv(csv_file_content: io.BytesIO):
    """Extracts text from a CSV file."""
    text = ""
    try:
        df = pd.read_csv(csv_file_content)
        text = df.to_string(index=False)
    except Exception as e:
        st.error(f"Error processing CSV: {e}")
    return text

def extract_text_from_html(html_file_content: io.BytesIO):
    """Extracts text from an HTML file."""
    text = ""
    try:
        soup = BeautifulSoup(html_file_content.read().decode("utf-8"), 'html.parser')
        text = soup.get_text(separator=' ', strip=True) # Extract visible text
    except Exception as e:
        st.error(f"Error processing HTML: {e}")
    return text

def extract_text_from_txt(txt_file_content: io.BytesIO):
    """Extracts text from a plain text (.txt) file."""
    text = ""
    try:
        text = txt_file_content.read().decode("utf-8")
    except Exception as e:
        st.error(f"Error processing TXT: {e}")
    return text

def extract_text(file_content: io.BytesIO, file_name: str):
    """
    Generic function to extract text based on file extension.
    It dispatches to specific extractors.
    """
    text = ""
    file_extension = file_name.split(".")[-1].lower()
    
    file_content.seek(0) # Reset stream position to the beginning before reading

    # Mapping of file extensions to their respective extraction functions
    extractors = {
        "pdf": extract_text_from_pdf,
        "pptx": extract_text_from_ppt,
        "py": extract_text_from_py,
        "doc": extract_text_from_docx, # handles both .doc and .docx
        "docx": extract_text_from_docx,
        "xls": extract_text_from_excel,
        "xlsx": extract_text_from_excel,
        "csv": extract_text_from_csv,
        "html": extract_text_from_html,
        "txt": extract_text_from_txt,
    }

    if file_extension in extractors:
        text = extractors[file_extension](file_content)
    else:
        st.warning(f"File type '.{file_extension}' not supported for text extraction. Skipping.")
    
    return text

# --- Document Processing and Vector Store Creation ---
@st.cache_resource(show_spinner="Processing documents and building knowledge base...")
def process_files_and_create_vectorstore(uploaded_files_list):
    """
    Processes a list of uploaded files, extracts text, splits it into chunks,
    and creates a FAISS vector store.
    """
    full_raw_text = ""
    doc_names = []

    # Iterate through each uploaded file to extract text
    for uploaded_file in uploaded_files_list:
        file_content = uploaded_file.getvalue()
        try:
            temp_bytes_io = io.BytesIO(file_content)
            extracted_content = extract_text(temp_bytes_io, uploaded_file.name)
            if extracted_content:
                full_raw_text += extracted_content + "\n\n--- Document Separator ---\n\n"
                doc_names.append(uploaded_file.name)
        except Exception as e:
            st.error(f"Could not process {uploaded_file.name}: {e}")

    # Check if any text was extracted
    if not full_raw_text.strip():
        st.warning("No readable content was extracted from the uploaded files.")
        st.session_state.uploaded_doc_names = []
        st.session_state.full_document_content = ""
        return None

    # Initialize text splitter
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    
    # Create document chunks. The metadata helps in identifying sources later.
    splits = text_splitter.create_documents(
        [full_raw_text],
        metadatas=[{"source": ", ".join(doc_names) if doc_names else "Uploaded Documents"}]
    )

    # Check if text splitting produced any chunks
    if not splits:
        st.warning("No text chunks could be generated from the documents.")
        st.session_state.uploaded_doc_names = []
        st.session_state.full_document_content = ""
        return None

    try:
        # Create embeddings and build the FAISS vector store
        embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
        vectorstore = FAISS.from_documents(splits, embeddings)
        
        # Store relevant information in session state for later use
        st.session_state.uploaded_doc_names = doc_names
        st.session_state.full_document_content = full_raw_text
        st.session_state.document_chunks = splits
        
        return vectorstore
    except Exception as e:
        st.error(f"Error creating vector store: {e}")
        return None
